"""upload_service — PDF/DOCX/PPTX/Markdown/TXT ingestion (Phase 27).

Extract-then-discard: parses the upload in-memory and calls the shared
`ingestion_service.ingest_document`, then never persists the original
bytes anywhere (no Supabase Storage bucket exists, and adding one is out
of scope for this phase — see the plan's locked decision).
"""

from __future__ import annotations

import io

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

from .document_repository import DocumentRepository
from .ingestion_service import ingest_document
from .km_models import KnowledgeDocument
from .source_repository import SourceRepository

SUPPORTED_UPLOAD_EXTENSIONS: set[str] = {"pdf", "docx", "pptx", "md", "markdown", "txt"}


def _extension(filename: str) -> str:
    return filename.lower().rsplit(".", 1)[-1] if "." in filename else ""


def extract_text(filename: str, content_bytes: bytes) -> str:
    ext = _extension(filename)

    if ext == "pdf":
        reader = PdfReader(io.BytesIO(content_bytes))
        return "\n\n".join(page.extract_text() or "" for page in reader.pages)

    if ext == "docx":
        document = DocxDocument(io.BytesIO(content_bytes))
        return "\n\n".join(p.text for p in document.paragraphs if p.text)

    if ext == "pptx":
        presentation = Presentation(io.BytesIO(content_bytes))
        texts = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame.text:
                    texts.append(shape.text_frame.text)
        return "\n\n".join(texts)

    if ext in ("md", "markdown", "txt"):
        return content_bytes.decode("utf-8", errors="ignore")

    raise ValueError(f"Unsupported file type: .{ext}")


class UploadService:
    def __init__(
        self,
        source_repository: SourceRepository | None = None,
        document_repository: DocumentRepository | None = None,
    ) -> None:
        self._sources = source_repository or SourceRepository()
        self._documents = document_repository or DocumentRepository()

    def ingest_upload(self, source_id: str, filename: str, content_bytes: bytes) -> KnowledgeDocument:
        source = self._sources.get(source_id)
        if source is None:
            raise ValueError(f"Unknown source: {source_id}")

        text = extract_text(filename, content_bytes)
        document = self._documents.create(source.workspace_id, source_id, parent_url=None, title=filename)

        # Extracted document text isn't HTML/markdown boilerplate the way
        # crawled pages are — skip the crawl-oriented cleaner.
        ingest_document(
            workspace_id=source.workspace_id,
            document_id=document.id,
            raw_text=text,
            parent_url=filename,
            product=source.product,
            source_type=source.source_type,
            clean=False,
        )
        result = self._documents.get(document.id)

        # ingest_document only ever updates the DOCUMENT's own status — the
        # crawl path (CrawlService) separately marks its SOURCE indexed/
        # ready once done, but nothing did that for uploads, leaving every
        # upload source stuck at "pending" forever even after a fully
        # successful ingest. Mirror CrawlService's success/failure handling
        # here so the source reflects what actually happened.
        if result is not None and result.status == "ready":
            self._sources.mark_indexed(source_id)
            self._sources.set_status(source_id, "ready")
        else:
            self._sources.set_status(source_id, "failed")

        return result
